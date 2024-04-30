from flask import Flask, request, jsonify
import os
from datetime import datetime
from PyPDF2 import PdfReader
from pptx import Presentation
import nltk
from nltk.tokenize import sent_tokenize, word_tokenize
from nltk import pos_tag
from random import choice
import matplotlib.pyplot as plt
from io import BytesIO
import base64
from PIL import Image

nltk.download('punkt')
nltk.download('averaged_perceptron_tagger')

app = Flask(__name__)

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files['file']
    original_filename = file.filename

    timestamp = datetime.now().strftime('%Y-%m-%d_%H-%M-%S')

    if original_filename.endswith('.pdf'):
        new_filename = f"{timestamp}_{original_filename}"
        file_path = os.path.join('data_pdf', new_filename)
        file.save(file_path)
        txt_file_path = convert_pdf_to_txt(file_path)
        preview_image = generate_pdf_preview(file_path)
    elif original_filename.endswith('.ppt') or original_filename.endswith('.pptx'):
        new_filename = f"{timestamp}_{original_filename}"
        file_path = os.path.join('data_ppt', new_filename)
        file.save(file_path)
        txt_file_path = convert_ppt_to_txt(file_path)
        preview_image = generate_ppt_preview(file_path)
    else:
        return jsonify({"error": "Unsupported file format"}), 400

    questions = generate_questions(txt_file_path)

    return jsonify({"message": "File processed successfully", "questions": questions, "preview_image": preview_image}), 200

def convert_pdf_to_txt(file_path):
    txt_file_path = file_path.replace('data_pdf/', 'data_txt/').replace('.pdf', '.txt')

    lines = []

    reader = PdfReader(file_path)
    for page in reader.pages:
        text = page.extract_text()
        lines.extend(text.split('\n'))

    with open(txt_file_path, 'w') as txt_file:
        current_paragraph = []

        for line in lines:
            line = line.strip()

            if not line:
                continue

            # Cek jika baris berakhir dengan titik
            if line.endswith('.'):
                current_paragraph.append(line)
                txt_file.write(" ".join(current_paragraph) + '\n\n')
                current_paragraph = []
            else:
                current_paragraph.append(line)

    return txt_file_path

def convert_ppt_to_txt(file_path):
    txt_file_path = file_path.replace('data_ppt/', 'data_txt/').replace('.pptx', '.txt').replace('.ppt', '.txt')

    prs = Presentation(file_path)
    with open(txt_file_path, 'w') as txt_file:
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    txt_file.write(shape.text + '\n')

    return txt_file_path

def generate_questions(txt_file_path, num_questions=3):
    # Baca teks dari file TXT
    with open(txt_file_path, 'r') as txt_file:
        text = txt_file.read()

    # Tokenisasi teks menjadi kalimat
    sentences = sent_tokenize(text)
    questions = []

    for _ in range(num_questions):
        # Pilih kalimat secara acak
        sentence = choice(sentences)

        # Tokenisasi kalimat menjadi kata-kata
        words = word_tokenize(sentence)

        # POS tagging untuk menemukan kata yang berarti
        tagged_words = pos_tag(words)

        # Pilih kata yang akan dihapus dari kata benda atau kata sifat
        valid_words = [w for w in tagged_words if w[1].startswith(('NN', 'JJ'))]

        if valid_words:
            word_to_remove = choice(valid_words)[0]

            # Bentuk pertanyaan fill-in-the-blank
            before, after = sentence.split(word_to_remove, 1)
            question = {
                "question": f"{before}______{after}",
                "before": before,
                "after": after,
                "answer": word_to_remove
            }

            questions.append(question)

    return questions

def generate_pdf_preview(file_path):
    from PyPDF2 import PdfReader
    from pdf2image import convert_from_path

    # Konversi halaman pertama menjadi gambar
    images = convert_from_path(file_path, first_page=1, last_page=1)

    # Simpan gambar ke dalam buffer
    img_buffer = BytesIO()
    images[0].save(img_buffer, format="PNG")

    # Encoding base64 untuk dikembalikan sebagai string
    encoded_image = base64.b64encode(img_buffer.getvalue()).decode('utf-8')

    return encoded_image

def generate_ppt_preview(file_path):
    prs = Presentation(file_path)

    # Ambil slide pertama
    first_slide = prs.slides[0]

    # Buat buffer untuk gambar
    slide_buffer = BytesIO()

    # Konversi slide menjadi gambar
    fig, ax = plt.subplots()
    ax.axis("off")

    for shape in first_slide.shapes:
        if hasattr(shape, "text"):
            ax.text(0.5, 0.5, shape.text, va="center", ha="center", wrap=True)

    plt.savefig(slide_buffer, format="PNG")
    slide_buffer.seek(0)

    # Encoding base64 untuk dikembalikan sebagai string
    encoded_image = base64.b64encode(slide_buffer.getvalue()).decode('utf-8')

    return encoded_image

if __name__ == '__main__':
    app.run(debug=True)
