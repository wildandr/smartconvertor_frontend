from flask import Flask, request, jsonify
import os
from datetime import datetime
from PyPDF2 import PdfReader
from pptx import Presentation
import nltk
from nltk.tokenize import sent_tokenize, word_tokenize
from nltk import pos_tag
from random import choice

nltk.download('punkt')
nltk.download('averaged_perceptron_tagger')

app = Flask(__name__)

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files['file']
    original_filename = file.filename

    timestamp = datetime.now().strftime('%Y-%m-%d_%H-%M-%S')

    if original_filename.endswith('.pdf'):
        new_filename = f"{timestamp}_{original_filename}"
        file_path = os.path.join('/Users/owwl/Documents/Joki/quiz_app/data_pdf', new_filename)
        file.save(file_path)
        txt_file_path = convert_pdf_to_txt(file_path)
    elif original_filename.endswith('.ppt') or original_filename.endswith('.pptx'):
        new_filename = f"{timestamp}_{original_filename}"
        file_path = os.path.join('/Users/owwl/Documents/Joki/quiz_app/data_ppt', new_filename)
        file.save(file_path)
        txt_file_path = convert_ppt_to_txt(file_path)
    else:
        return jsonify({"error": "Unsupported file format"}), 400

    questions = generate_questions(txt_file_path)

    return jsonify({"message": "File processed successfully", "questions": questions}), 200

def convert_pdf_to_txt(file_path):
    txt_file_path = file_path.replace('/data_pdf/', '/data_txt/').replace('.pdf', '.txt')

    lines = []

    reader = PdfReader(file_path)
    for page in reader.pages:
        text = page.extract_text()
        lines.extend(text.split('\n'))

    with open(txt_file_path, 'w') as txt_file:
        current_paragraph = []

        for line in lines:
            line = line.strip()

            if not line:
                continue

            # Cek jika baris berakhir dengan titik
            if line.endswith('.'):
                current_paragraph.append(line)
                txt_file.write(" ".join(current_paragraph) + '\n\n')
                current_paragraph = []
            else:
                current_paragraph.append(line)

    return txt_file_path

def convert_ppt_to_txt(file_path):
    txt_file_path = file_path.replace('/data_ppt/', '/data_txt/').replace('.pptx', '.txt').replace('.ppt', '.txt')

    prs = Presentation(file_path)
    with open(txt_file_path, 'w') as txt_file:
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    txt_file.write(shape.text + '\n')

    return txt_file_path

def generate_questions(txt_file_path, num_questions=3):
    # Baca teks dari file TXT
    with open(txt_file_path, 'r') as txt_file:
        text = txt_file.read()

    # Tokenisasi teks menjadi kalimat
    sentences = sent_tokenize(text)
    questions = []

    for _ in range(num_questions):
        # Pilih kalimat secara acak
        sentence = choice(sentences)

        # Tokenisasi kalimat menjadi kata-kata
        words = word_tokenize(sentence)

        # POS tagging untuk menemukan kata yang berarti
        tagged_words = pos_tag(words)

        # Pilih kata yang akan dihapus dari kata benda atau kata sifat
        valid_words = [w for w in tagged_words if w[1].startswith(('NN', 'JJ'))]

        if valid_words:
            word_to_remove = choice(valid_words)[0]

            # Bentuk pertanyaan fill-in-the-blank
            before, after = sentence.split(word_to_remove, 1)
            question = {
                "question": f"{before}______{after}",
                "before": before,
                "after": after,
                "answer": word_to_remove
            }

            questions.append(question)

    return questions

if __name__ == '__main__':
    app.run(debug=True)
